from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "PassportSOL_Pitch_Collector_v2.pptx"

IMG_HERO = ROOT / "src" / "assets" / "hero.png"
IMG_FLOW = ROOT / "slide-3.jpg"
IMG_BADGES = ROOT / "slide-4.jpg"
IMG_PASSPORT = ROOT / "slide-5.jpg"
IMG_WIDGET = ROOT / "slide-7.jpg"


BG = RGBColor(246, 240, 231)
PANEL = RGBColor(255, 251, 245)
INK = RGBColor(27, 27, 24)
MUTED = RGBColor(109, 105, 98)
FOREST = RGBColor(25, 58, 46)
MINT = RGBColor(45, 175, 122)
GOLD = RGBColor(180, 131, 59)
LINE = RGBColor(214, 202, 187)
WHITE = RGBColor(255, 255, 255)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_round_box(slide, left, top, width, height, fill_color, line_color=None, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color or fill_color
    shape.line.width = Pt(1)
    return shape


def add_textbox(slide, left, top, width, height, text="", font_name="Aptos", font_size=18,
                color=INK, bold=False, italic=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    font = r.font
    font.name = font_name
    font.size = Pt(font_size)
    font.color.rgb = color
    font.bold = bold
    font.italic = italic
    return box


def add_bullets(slide, left, top, width, height, items, font_size=20, color=INK, bullet_color=FOREST):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for i, item in enumerate(items):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        r1 = p.add_run()
        r1.text = "• "
        r1.font.name = "Aptos"
        r1.font.size = Pt(font_size)
        r1.font.color.rgb = bullet_color
        r1.font.bold = True
        r2 = p.add_run()
        r2.text = item
        r2.font.name = "Aptos"
        r2.font.size = Pt(font_size)
        r2.font.color.rgb = color
    return box


def add_kicker(slide, text):
    add_textbox(slide, Inches(0.75), Inches(0.45), Inches(3.8), Inches(0.35), text,
                font_name="Aptos", font_size=12, color=FOREST, bold=True)


def add_title(slide, text, left=0.75, top=0.9, width=6.0, height=1.2, size=28):
    add_textbox(slide, Inches(left), Inches(top), Inches(width), Inches(height), text,
                font_name="Georgia", font_size=size, color=INK, bold=True)


def add_image_panel(slide, img_path, left, top, width, height):
    frame = add_round_box(slide, left, top, width, height, PANEL, LINE)
    frame.shadow.inherit = False
    slide.shapes.add_picture(str(img_path), left + Inches(0.08), top + Inches(0.08),
                             width=width - Inches(0.16), height=height - Inches(0.16))
    return frame


def add_chip(slide, left, top, text, fill_color=PANEL, text_color=FOREST, width=None):
    width = width or Inches(1.5)
    chip = add_round_box(slide, left, top, width, Inches(0.34), fill_color, LINE)
    tf = chip.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(11)
    r.font.color.rgb = text_color
    r.font.bold = True
    return chip


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_round_box(slide, Inches(0.7), Inches(0.55), Inches(1.8), Inches(0.28), PANEL, LINE)
    add_textbox(slide, Inches(0.83), Inches(0.60), Inches(1.5), Inches(0.18), "COLLECTIBLE LOYALTY",
                font_size=10, color=FOREST, bold=True)
    add_title(slide, "PassportSOL", top=1.05, size=30)
    add_textbox(
        slide, Inches(0.75), Inches(1.8), Inches(5.7), Inches(1.8),
        "Collect badges. Build reputation. Unlock rewards.",
        font_name="Georgia", font_size=26, color=FOREST, bold=True
    )
    add_textbox(
        slide, Inches(0.75), Inches(3.2), Inches(5.2), Inches(1.2),
        "PassportSOL turns on-chain identity into a collectible reward loop that projects can trust and users actually want.",
        font_size=18, color=MUTED
    )
    add_chip(slide, Inches(0.78), Inches(4.5), "Free for users", width=Inches(1.4))
    add_chip(slide, Inches(2.28), Inches(4.5), "Built for Solana", width=Inches(1.6))
    add_chip(slide, Inches(4.02), Inches(4.5), "Sybil resistant", width=Inches(1.7))
    add_round_box(slide, Inches(7.0), Inches(0.8), Inches(5.0), Inches(5.6), WHITE, LINE)
    slide.shapes.add_picture(str(IMG_HERO), Inches(7.85), Inches(1.2), height=Inches(4.8))
    add_textbox(slide, Inches(7.45), Inches(6.0), Inches(4.2), Inches(0.6),
                "A loyalty passport for communities that reward real participation.",
                font_size=15, color=MUTED, align=PP_ALIGN.CENTER)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_kicker(slide, "THE SHIFT")
    add_title(slide, "Airdrops reward speed. Communities want to reward belonging.", top=0.95, width=8.2, size=24)
    stat = add_round_box(slide, Inches(0.8), Inches(1.85), Inches(3.0), Inches(2.0), FOREST, FOREST)
    tf = stat.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "$94.5M"
    r.font.name = "Georgia"
    r.font.size = Pt(28)
    r.font.color.rgb = WHITE
    r.font.bold = True
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "drained by sybil bots in one airdrop"
    r2.font.name = "Aptos"
    r2.font.size = Pt(13)
    r2.font.color.rgb = WHITE
    add_bullets(
        slide, Inches(4.25), Inches(1.9), Inches(7.6), Inches(2.4),
        [
            "Projects struggle to tell a loyal fan from a temporary wallet farm.",
            "Users get one more claim page instead of something fun to collect and keep.",
            "Cross-chain users still need an easy first step into Solana before they can participate.",
        ],
        font_size=18
    )
    add_round_box(slide, Inches(0.8), Inches(4.35), Inches(11.7), Inches(1.65), PANEL, LINE)
    add_textbox(slide, Inches(1.0), Inches(4.72), Inches(11.2), Inches(0.85),
                "The real opportunity is not another whitelist. It is a collectible loyalty layer that makes participation visible, portable, and worth building over time.",
                font_name="Georgia", font_size=20, color=INK)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_kicker(slide, "HOW IT WORKS")
    add_title(slide, "Simple for users. Sticky for communities.", top=0.95, width=6.0, size=24)
    add_bullets(
        slide, Inches(0.85), Inches(1.95), Inches(4.6), Inches(3.8),
        [
            "Fund a Solana wallet in one smooth step, even if the user starts on another chain.",
            "Collect badges for real signals like events, dev identity, and community presence.",
            "Mint one passport that turns those signals into a reusable reputation card.",
            "Unlock rewards that feel earned, not sprayed at random wallets.",
        ],
        font_size=17
    )
    add_chip(slide, Inches(0.9), Inches(5.7), "Smooth entry", width=Inches(1.5))
    add_chip(slide, Inches(2.55), Inches(5.7), "Collectible loop", width=Inches(1.8))
    add_chip(slide, Inches(4.5), Inches(5.7), "Portable proof", width=Inches(1.7))
    add_image_panel(slide, IMG_FLOW, Inches(6.15), Inches(1.15), Inches(6.2), Inches(5.65))

    # Slide 4
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_kicker(slide, "COLLECTOR MOMENT")
    add_title(slide, "Badges make participation visible.", top=0.95, width=5.4, size=24)
    add_textbox(slide, Inches(0.85), Inches(1.9), Inches(4.2), Inches(1.9),
                "Every badge gives the user something to show: attended, built, claimed, belonged. That is what turns a one-time campaign into a habit.",
                font_size=18, color=MUTED)
    add_bullets(
        slide, Inches(0.9), Inches(3.7), Inches(4.2), Inches(2.3),
        [
            "Feels like a loyalty card, not a compliance form.",
            "Encourages repeat engagement because people want the full set.",
            "Raises reward quality because badges are harder for bots to fake.",
        ],
        font_size=17
    )
    add_image_panel(slide, IMG_BADGES, Inches(5.45), Inches(1.1), Inches(6.65), Inches(5.95))

    # Slide 5
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_kicker(slide, "THE PASSPORT")
    add_title(slide, "One card. All your proof.", top=0.95, width=5.3, size=24)
    add_image_panel(slide, IMG_PASSPORT, Inches(0.82), Inches(1.35), Inches(6.15), Inches(5.5))
    add_round_box(slide, Inches(7.35), Inches(1.5), Inches(4.8), Inches(5.0), WHITE, LINE)
    add_textbox(slide, Inches(7.7), Inches(1.95), Inches(4.0), Inches(0.7),
                "A passport people actually want to keep:", font_size=18, color=FOREST, bold=True)
    add_bullets(
        slide, Inches(7.7), Inches(2.55), Inches(3.95), Inches(2.8),
        [
            "Collectible: your progress is visible and shareable.",
            "Portable: one identity can travel across Solana apps and communities.",
            "Compounding: the more you collect, the stronger your reputation becomes.",
        ],
        font_size=17
    )
    add_textbox(slide, Inches(7.7), Inches(5.55), Inches(3.95), Inches(0.85),
                "Projects do not need to rebuild trust from scratch every time. Users bring their passport with them.",
                font_size=15, color=MUTED)

    # Slide 6
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_kicker(slide, "WHY IT WINS")
    add_title(slide, "PassportSOL makes sybil defense feel invisible.", top=0.95, width=6.2, size=24)
    add_image_panel(slide, IMG_WIDGET, Inches(0.82), Inches(1.35), Inches(6.25), Inches(5.5))
    add_bullets(
        slide, Inches(7.45), Inches(1.7), Inches(4.2), Inches(3.2),
        [
            "For users: a smooth entry point, clear progress, and rewards that feel earned.",
            "For projects: fewer fake claimants and a better way to find real community members.",
            "For the product: security sits in the background while the front-end experience stays playful and collectible.",
        ],
        font_size=17
    )
    add_round_box(slide, Inches(7.45), Inches(5.15), Inches(4.15), Inches(1.05), PANEL, LINE)
    add_textbox(slide, Inches(7.7), Inches(5.45), Inches(3.7), Inches(0.5),
                "The experience leads with collecting. The defense is built in.",
                font_size=16, color=INK)

    # Slide 7
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_kicker(slide, "BUSINESS MODEL")
    add_title(slide, "Projects pay. Users never do.", top=0.95, width=5.8, size=24)
    add_round_box(slide, Inches(0.85), Inches(1.8), Inches(2.25), Inches(1.2), FOREST, FOREST)
    add_textbox(slide, Inches(1.1), Inches(2.12), Inches(1.75), Inches(0.4), "$0 user cost",
                font_name="Georgia", font_size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_round_box(slide, Inches(0.85), Inches(3.25), Inches(5.3), Inches(2.45), WHITE, LINE)
    add_textbox(slide, Inches(1.15), Inches(3.6), Inches(2.1), Inches(0.45), "Starter",
                font_name="Georgia", font_size=24, color=FOREST, bold=True)
    add_textbox(slide, Inches(1.15), Inches(4.05), Inches(2.3), Inches(0.55), "$49 / month",
                font_name="Georgia", font_size=28, color=INK, bold=True)
    add_textbox(slide, Inches(1.15), Inches(4.72), Inches(4.35), Inches(0.7),
                "A fixed monthly price for smaller launches, new communities, and pilot campaigns.",
                font_size=16, color=MUTED)
    add_round_box(slide, Inches(6.45), Inches(3.25), Inches(5.45), Inches(2.45), PANEL, LINE)
    add_textbox(slide, Inches(6.75), Inches(3.6), Inches(2.7), Inches(0.45), "Larger projects",
                font_name="Georgia", font_size=24, color=FOREST, bold=True)
    add_textbox(slide, Inches(6.75), Inches(4.05), Inches(3.0), Inches(0.55), "Pay as you go",
                font_name="Georgia", font_size=28, color=INK, bold=True)
    add_textbox(slide, Inches(6.75), Inches(4.72), Inches(4.45), Inches(0.7),
                "Volume-based pricing for bigger drops, bigger communities, and higher verification demand.",
                font_size=16, color=MUTED)
    add_chip(slide, Inches(0.95), Inches(6.2), "Memecoin launches", width=Inches(1.7))
    add_chip(slide, Inches(2.8), Inches(6.2), "NFT communities", width=Inches(1.65))
    add_chip(slide, Inches(4.55), Inches(6.2), "DAOs", width=Inches(0.85))
    add_chip(slide, Inches(5.55), Inches(6.2), "Games", width=Inches(0.95))
    add_chip(slide, Inches(6.65), Inches(6.2), "DeFi campaigns", width=Inches(1.5))

    # Slide 8
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_round_box(slide, Inches(0.75), Inches(0.85), Inches(11.75), Inches(5.7), WHITE, LINE)
    slide.shapes.add_picture(str(IMG_HERO), Inches(0.95), Inches(1.25), height=Inches(4.65))
    add_textbox(slide, Inches(4.65), Inches(1.4), Inches(6.8), Inches(0.4), "PASSPORTSOL",
                font_size=12, color=FOREST, bold=True)
    add_textbox(slide, Inches(4.65), Inches(1.85), Inches(6.6), Inches(1.2),
                "The collectible reputation layer for Solana communities.",
                font_name="Georgia", font_size=28, color=INK, bold=True)
    add_textbox(slide, Inches(4.65), Inches(3.25), Inches(6.4), Inches(1.2),
                "Users collect badges. Communities reward real participation. Bots have less room to hide.",
                font_size=18, color=MUTED)
    add_textbox(slide, Inches(4.65), Inches(4.7), Inches(5.8), Inches(0.8),
                "The more you collect, the more you unlock.",
                font_name="Georgia", font_size=24, color=FOREST, bold=True)
    add_textbox(slide, Inches(0.85), Inches(6.85), Inches(11.4), Inches(0.3),
                "Repo: github.com/0jz/passportsol", font_size=11, color=MUTED, align=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f"Created: {OUTPUT}")


if __name__ == "__main__":
    build_deck()
