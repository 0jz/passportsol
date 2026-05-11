from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "PassportSOL_Pitch_Graphics_Short_v3.pptx"


BG = RGBColor(248, 244, 237)
CARD = RGBColor(255, 252, 247)
INK = RGBColor(24, 24, 29)
MUTED = RGBColor(108, 104, 100)
FOREST = RGBColor(32, 74, 58)
PLUM = RGBColor(99, 49, 194)
TEAL = RGBColor(26, 176, 140)
CORAL = RGBColor(245, 122, 93)
GOLD = RGBColor(199, 143, 68)
LINE = RGBColor(214, 205, 194)
WHITE = RGBColor(255, 255, 255)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_shape(slide, shape_type, left, top, width, height, fill, line=None, line_width=1):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(line_width)
    return shape


def add_round(slide, left, top, width, height, fill, line=None, line_width=1):
    return add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
        fill,
        line,
        line_width,
    )


def add_text(slide, left, top, width, height, text, *, font="Aptos", size=18, color=INK,
             bold=False, italic=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return box


def add_bullets(slide, left, top, width, height, items, *, size=18):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        p.alignment = PP_ALIGN.LEFT
        a = p.add_run()
        a.text = "• "
        a.font.name = "Aptos"
        a.font.size = Pt(size)
        a.font.color.rgb = PLUM
        a.font.bold = True
        b = p.add_run()
        b.text = item
        b.font.name = "Aptos"
        b.font.size = Pt(size)
        b.font.color.rgb = INK
    return box


def add_kicker(slide, text):
    add_text(slide, Inches(0.8), Inches(0.45), Inches(3.5), Inches(0.3), text,
             size=11, color=PLUM, bold=True)


def add_title(slide, text, *, top=0.85, width=7.0, size=28):
    add_text(slide, Inches(0.8), Inches(top), Inches(width), Inches(0.9), text,
             font="Georgia", size=size, color=INK, bold=True)


def add_gradient_like_bars(slide):
    add_round(slide, Inches(10.7), Inches(0.5), Inches(1.8), Inches(0.22), PLUM, PLUM, 0)
    add_round(slide, Inches(10.0), Inches(0.82), Inches(2.5), Inches(0.18), TEAL, TEAL, 0)
    add_round(slide, Inches(10.9), Inches(1.08), Inches(1.6), Inches(0.14), CORAL, CORAL, 0)


def add_graphic_cluster(slide, left, top, scale=1.0):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, left, top, Inches(1.45 * scale), Inches(1.45 * scale), CARD, LINE)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.95 * scale), top + Inches(0.3 * scale),
              Inches(1.15 * scale), Inches(1.15 * scale), TEAL, TEAL)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(1.8 * scale), top + Inches(0.95 * scale),
              Inches(0.8 * scale), Inches(0.8 * scale), GOLD, GOLD)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.DIAMOND, left + Inches(0.45 * scale), top + Inches(1.25 * scale),
              Inches(0.6 * scale), Inches(0.6 * scale), CORAL, CORAL)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, left + Inches(1.55 * scale), top - Inches(0.08 * scale),
              Inches(0.55 * scale), Inches(0.55 * scale), PLUM, PLUM)


def add_screenshot_placeholder(slide, left, top, width, height, label):
    frame = add_round(slide, left, top, width, height, CARD, LINE)
    add_round(slide, left + Inches(0.16), top + Inches(0.18), width - Inches(0.32), Inches(0.28), INK, INK, 0)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.3), top + Inches(0.25), Inches(0.08), Inches(0.08), CORAL, CORAL, 0)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.42), top + Inches(0.25), Inches(0.08), Inches(0.08), GOLD, GOLD, 0)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.54), top + Inches(0.25), Inches(0.08), Inches(0.08), TEAL, TEAL, 0)
    add_text(slide, left + Inches(0.45), top + Inches(1.4), width - Inches(0.9), Inches(0.5), label,
             size=16, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.45), top + Inches(1.85), width - Inches(0.9), Inches(0.45), "drop app screenshot here",
             size=12, color=MUTED, align=PP_ALIGN.CENTER)
    return frame


def add_metric_chip(slide, left, top, title, value, accent):
    chip = add_round(slide, left, top, Inches(2.25), Inches(1.05), WHITE, LINE)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.18), Inches(1.05), accent, accent, 0)
    add_text(slide, left + Inches(0.32), top + Inches(0.18), Inches(1.5), Inches(0.24), title,
             size=10, color=MUTED, bold=True)
    add_text(slide, left + Inches(0.32), top + Inches(0.42), Inches(1.7), Inches(0.35), value,
             font="Georgia", size=22, color=INK, bold=True)
    return chip


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 cover
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_gradient_like_bars(slide)
    add_text(slide, Inches(0.82), Inches(0.52), Inches(2.8), Inches(0.28), "NEW DECK / SHORT VERSION",
             size=10, color=MUTED, bold=True)
    add_title(slide, "PassportSOL", top=1.0, width=5.4, size=31)
    add_text(slide, Inches(0.82), Inches(1.9), Inches(5.4), Inches(1.1),
             "Collect badges.\nUnlock rewards.",
             font="Georgia", size=24, color=PLUM, bold=True)
    add_text(slide, Inches(0.82), Inches(3.2), Inches(5.3), Inches(0.75),
             "A collectible loyalty layer for Solana communities.",
             size=18, color=MUTED)
    add_graphic_cluster(slide, Inches(7.8), Inches(1.15), 1.35)
    add_round(slide, Inches(7.0), Inches(4.45), Inches(5.1), Inches(1.15), WHITE, LINE)
    add_text(slide, Inches(7.35), Inches(4.78), Inches(4.4), Inches(0.45),
             "Short story, bold visuals, screenshot-ready layout.",
             size=16, color=INK)

    # 2 problem
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_kicker(slide, "PROBLEM")
    add_title(slide, "Airdrops feel random.", top=0.9, width=4.8, size=28)
    add_text(slide, Inches(0.82), Inches(1.75), Inches(3.1), Inches(0.8),
             "Real fans get treated like anonymous wallets.", size=19, color=MUTED)
    add_metric_chip(slide, Inches(0.85), Inches(3.0), "bot loss example", "$94.5M", CORAL)
    add_metric_chip(slide, Inches(3.3), Inches(3.0), "user signal", "too weak", GOLD)
    add_metric_chip(slide, Inches(5.75), Inches(3.0), "collector value", "missing", TEAL)
    add_round(slide, Inches(8.75), Inches(1.4), Inches(3.2), Inches(4.8), WHITE, LINE)
    add_text(slide, Inches(9.1), Inches(1.9), Inches(2.6), Inches(0.55), "What’s missing?", size=20, color=INK, bold=True)
    add_bullets(slide, Inches(9.08), Inches(2.55), Inches(2.55), Inches(2.7),
                ["proof of loyalty", "something collectible", "clean reward gating"], size=16)

    # 3 product
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_kicker(slide, "PRODUCT")
    add_title(slide, "PassportSOL makes loyalty visible.", top=0.9, width=6.4, size=28)
    add_bullets(slide, Inches(0.85), Inches(1.85), Inches(4.7), Inches(2.9),
                ["Collect badges", "Build one passport", "Use it across communities"], size=19)
    add_round(slide, Inches(0.88), Inches(5.15), Inches(4.9), Inches(1.0), WHITE, LINE)
    add_text(slide, Inches(1.15), Inches(5.47), Inches(4.3), Inches(0.36),
             "Short version: it feels more like a loyalty game than a claim form.",
             size=15, color=MUTED)
    add_screenshot_placeholder(slide, Inches(6.55), Inches(1.25), Inches(5.7), Inches(5.4), "APP SHOT 01")

    # 4 badges
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_kicker(slide, "COLLECT")
    add_title(slide, "Badges are the hook.", top=0.9, width=4.8, size=28)
    add_round(slide, Inches(0.85), Inches(1.8), Inches(2.15), Inches(1.0), WHITE, LINE)
    add_text(slide, Inches(1.15), Inches(2.13), Inches(1.55), Inches(0.3), "events", size=18, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_round(slide, Inches(3.2), Inches(1.8), Inches(2.15), Inches(1.0), WHITE, LINE)
    add_text(slide, Inches(3.5), Inches(2.13), Inches(1.55), Inches(0.3), "dev proof", size=18, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_round(slide, Inches(5.55), Inches(1.8), Inches(2.15), Inches(1.0), WHITE, LINE)
    add_text(slide, Inches(5.85), Inches(2.13), Inches(1.55), Inches(0.3), ".sol / ENS", size=18, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_round(slide, Inches(0.88), Inches(3.35), Inches(6.8), Inches(1.2), CARD, LINE)
    add_text(slide, Inches(1.15), Inches(3.72), Inches(6.2), Inches(0.5),
             "More badges = stronger reputation = better rewards.",
             font="Georgia", size=23, color=FOREST, bold=True)
    add_screenshot_placeholder(slide, Inches(8.2), Inches(1.25), Inches(4.0), Inches(5.5), "APP SHOT 02")

    # 5 passport
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_kicker(slide, "PASSPORT")
    add_title(slide, "One card.\nAll your proof.", top=0.9, width=3.4, size=27)
    add_text(slide, Inches(0.86), Inches(2.4), Inches(3.5), Inches(1.0),
             "Portable, simple, and easy to verify.", size=18, color=MUTED)
    add_bullets(slide, Inches(0.9), Inches(3.5), Inches(3.8), Inches(2.2),
                ["shareable", "reusable", "cross-community"], size=18)
    add_screenshot_placeholder(slide, Inches(5.15), Inches(1.0), Inches(6.4), Inches(5.9), "APP SHOT 03")

    # 6 business model
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_kicker(slide, "BUSINESS MODEL")
    add_title(slide, "Simple pricing.", top=0.9, width=4.2, size=28)
    add_round(slide, Inches(0.88), Inches(1.9), Inches(4.9), Inches(3.15), WHITE, LINE)
    add_text(slide, Inches(1.2), Inches(2.25), Inches(1.7), Inches(0.35), "Starter", size=22, color=FOREST, bold=True)
    add_text(slide, Inches(1.2), Inches(2.65), Inches(2.4), Inches(0.5), "$49 / month", font="Georgia", size=28, color=INK, bold=True)
    add_text(slide, Inches(1.2), Inches(3.35), Inches(3.9), Inches(0.75), "Fixed monthly pricing for smaller launches and pilot communities.", size=16, color=MUTED)
    add_round(slide, Inches(6.15), Inches(1.9), Inches(5.1), Inches(3.15), CARD, LINE)
    add_text(slide, Inches(6.48), Inches(2.25), Inches(2.4), Inches(0.35), "Bigger projects", size=22, color=FOREST, bold=True)
    add_text(slide, Inches(6.48), Inches(2.65), Inches(2.7), Inches(0.5), "Pay as you go", font="Georgia", size=28, color=INK, bold=True)
    add_text(slide, Inches(6.48), Inches(3.35), Inches(4.1), Inches(0.75), "Usage-based pricing for larger campaigns and higher verification volume.", size=16, color=MUTED)
    add_round(slide, Inches(1.0), Inches(5.55), Inches(10.4), Inches(0.9), WHITE, LINE)
    add_text(slide, Inches(1.35), Inches(5.83), Inches(9.7), Inches(0.35),
             "Users never pay. Communities pay to reward real participation better.",
             size=16, color=INK, align=PP_ALIGN.CENTER)

    # 7 close
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_graphic_cluster(slide, Inches(0.95), Inches(1.2), 1.1)
    add_title(slide, "PassportSOL", top=1.0, width=4.0, size=31)
    add_text(slide, Inches(0.85), Inches(2.0), Inches(4.5), Inches(1.2),
             "Collect badges.\nUnlock rewards.",
             font="Georgia", size=25, color=PLUM, bold=True)
    add_text(slide, Inches(0.85), Inches(3.55), Inches(4.6), Inches(0.8),
             "A short, collectible pitch for a more human kind of reward system.",
             size=18, color=MUTED)
    add_screenshot_placeholder(slide, Inches(6.2), Inches(1.05), Inches(5.7), Inches(5.9), "HERO APP SHOT")
    add_text(slide, Inches(0.85), Inches(6.7), Inches(11.5), Inches(0.3),
             "github.com/0jz/passportsol", size=11, color=MUTED, align=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f"Created: {OUTPUT}")


if __name__ == "__main__":
    build()
