from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "PassportSOL_Pitch_Judges_v1.pptx"
SCRIPT_OUTPUT = ROOT / "docs" / "PASSPORTSOL_JUDGE_PITCH.md"


BG = RGBColor(245, 241, 234)
CARD = RGBColor(255, 252, 246)
INK = RGBColor(23, 25, 28)
MUTED = RGBColor(101, 101, 101)
FOREST = RGBColor(23, 73, 58)
TEAL = RGBColor(24, 165, 140)
CORAL = RGBColor(233, 107, 78)
GOLD = RGBColor(190, 146, 63)
LINE = RGBColor(215, 206, 194)
WHITE = RGBColor(255, 255, 255)


SLIDES = [
    {
        "kicker": "PASSPORTSOL",
        "title": "Gitcoin gives the seed.\nSolana grows the passport.",
        "body": [
            "PassportSOL turns Gitcoin trust plus public Solana activity into collectible stamps.",
            "Communities use that passport to reward real participation, not anonymous wallets.",
        ],
        "accent": "seed + signal",
        "visual": "signal map",
    },
    {
        "kicker": "PROBLEM",
        "title": "Airdrops still treat everyone like a fresh wallet.",
        "body": [
            "Gitcoin alone does not tell a Solana community who actually stayed, built, voted, collected, or showed up.",
            "That is why rewards leak to low-context wallets while real members look invisible.",
        ],
        "accent": "$94.5M lost",
        "visual": "wallet fog",
    },
    {
        "kicker": "INSIGHT",
        "title": "Reputation should feel earned, visible, and native to Solana.",
        "body": [
            "We start with Gitcoin as a cross-chain trust seed, then add stamps from open Solana data like wallet age, SNS, governance, collection history, and ecosystem activity.",
            "The result is a living passport that gets stronger as your on-chain story grows.",
        ],
        "accent": "living passport",
        "visual": "stamp ladder",
    },
    {
        "kicker": "PRODUCT",
        "title": "Collectible stamps are the product.",
        "body": [
            "Users do not just pass a check; they collect identity objects like Builder, Voter, Collector, Event Attendee, and Early Community.",
            "Projects can read the same passport to gate rewards, campaigns, and access with one shared reputation layer.",
        ],
        "accent": "collect, not just verify",
        "visual": "stamp grid",
    },
    {
        "kicker": "WHY SOLANA",
        "title": "Solana makes this more useful, not just more on-chain.",
        "body": [
            "Public wallet behavior, token activity, community signals, and low-cost updates make Solana the best place to turn reputation into a live product.",
            "Gitcoin gives credibility at the edge, but Solana gives us the ongoing data that makes the passport worth checking.",
        ],
        "accent": "native public data",
        "visual": "data rails",
    },
    {
        "kicker": "BUSINESS MODEL",
        "title": "Users collect for free.\nCommunities pay to use the passport.",
        "body": [
            "Starter is a fixed $49 per month for smaller launches and pilot communities.",
            "Larger projects pay as they go when they need more verification, segmentation, and reward gating.",
        ],
        "accent": "$49 + usage",
        "visual": "pricing split",
    },
    {
        "kicker": "CLOSE",
        "title": "PassportSOL makes Gitcoin useful inside Solana-native communities.",
        "body": [
            "Our wedge is simple: collectible stamps powered by Gitcoin trust and public Solana data.",
            "The more signal a wallet earns, the more confidently a community can reward it.",
        ],
        "accent": "reward real members",
        "visual": "hero shot",
    },
]


SCRIPT_TEXT = """# PassportSOL Judge Pitch

## One-liner
PassportSOL turns Gitcoin trust plus public Solana activity into collectible stamps that communities can reward.

## Two-minute pitch

Airdrops still treat everyone like a fresh wallet. Gitcoin can tell you a user has some cross-chain trust, but it still does not tell a Solana community who actually stayed, built, voted, collected, or showed up.

That is the gap PassportSOL fills.

We use Gitcoin as the starting trust seed, then grow that identity with public Solana data. Instead of a static score, users collect stamps from visible signals like wallet age, SNS ownership, governance participation, collection history, ecosystem activity, and event presence.

That changes the product from a gate into a collectible reputation system.

For users, the value is simple: they are not just verifying once, they are building a passport that gets richer over time. For communities, the value is also simple: they can finally reward wallets with context, not just wallets that arrived first.

This is why Solana matters here. The chain gives us rich public data, low-cost updates, and a natural environment for a living identity product. Gitcoin gives us credibility at the edge, but Solana gives us the ongoing signal that makes the passport worth checking.

Business-model wise, users collect for free. Smaller communities can start at a fixed $49 per month, and larger projects pay as they go when they need more verification, segmentation, and reward gating.

PassportSOL makes Gitcoin useful inside Solana-native communities by turning trust into something visible, collectible, and actionable.
"""


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_shape(slide, shape_type, left, top, width, height, fill_color, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color or fill_color
    shape.line.width = Pt(line_width)
    return shape


def add_round(slide, left, top, width, height, fill_color, line_color=None):
    return add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
        fill_color,
        line_color,
    )


def add_text(slide, left, top, width, height, text, *, font="Aptos", size=18, color=INK,
             bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return box


def add_body(slide, left, top, width, height, lines):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.name = "Aptos"
        run.font.size = Pt(19)
        run.font.color.rgb = MUTED if idx == 0 else INK
    return box


def add_header(slide, kicker, title):
    add_text(slide, Inches(0.78), Inches(0.5), Inches(2.8), Inches(0.28), kicker,
             size=10, color=FOREST, bold=True)
    add_text(slide, Inches(0.78), Inches(1.0), Inches(6.2), Inches(1.35), title,
             font="Georgia", size=28, color=INK, bold=True)


def add_accent_tag(slide, text):
    tag = add_round(slide, Inches(0.82), Inches(5.95), Inches(2.35), Inches(0.52), WHITE, LINE)
    tf = tag.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(12)
    run.font.color.rgb = FOREST
    run.font.bold = True


def add_visual_frame(slide, label):
    frame = add_round(slide, Inches(7.05), Inches(1.0), Inches(5.35), Inches(5.85), CARD, LINE)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, Inches(7.55), Inches(1.45), Inches(1.6), Inches(1.6), TEAL, TEAL)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.0), Inches(2.05), Inches(1.15), Inches(1.15), GOLD, GOLD)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.DIAMOND, Inches(10.05), Inches(1.55), Inches(0.9), Inches(0.9), CORAL, CORAL)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, Inches(8.0), Inches(3.1), Inches(1.0), Inches(1.0), FOREST, FOREST)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.HEXAGON, Inches(9.5), Inches(3.35), Inches(1.2), Inches(1.2), WHITE, LINE)
    add_text(slide, Inches(7.65), Inches(4.95), Inches(4.2), Inches(0.45), label.upper(),
             size=18, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(7.6), Inches(5.38), Inches(4.25), Inches(0.45), "native PPTX graphic block",
             size=12, color=MUTED, italic=True, align=PP_ALIGN.CENTER)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for slide_data in SLIDES:
        slide = prs.slides.add_slide(blank)
        set_bg(slide)
        add_header(slide, slide_data["kicker"], slide_data["title"])
        add_body(slide, Inches(0.82), Inches(2.55), Inches(5.65), Inches(2.7), slide_data["body"])
        add_accent_tag(slide, slide_data["accent"])
        add_visual_frame(slide, slide_data["visual"])

    prs.save(OUTPUT)
    SCRIPT_OUTPUT.write_text(SCRIPT_TEXT, encoding="utf-8")
    print(f"Created: {OUTPUT}")
    print(f"Created: {SCRIPT_OUTPUT}")


if __name__ == "__main__":
    build_deck()
